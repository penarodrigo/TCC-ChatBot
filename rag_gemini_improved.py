# -*- coding: utf-8 -*-
"""
Sistema RAG com Google Gemini para VSCode
Versão Revisada e Melhorada
Adaptado do Workshop Energy GPT - RAG
Focado em responder com base no documento IDDC.pdf
"""

# Importações de bibliotecas padrão do Python
import os
import logging
import pickle
import hashlib
from pathlib import Path
from typing import List, Dict, Optional, Union, Tuple
from datetime import datetime
import traceback # Para rastreamento de erros mais detalhado, se necessário
import json # Para manipulação de dados JSON
import argparse # Para processar argumentos de linha de comando
import sys # Para sys.exit e sys.argv

# Importações para processamento de documentos
import pandas as pd
import pdfplumber
from pptx import Presentation
import docx # Para arquivos .docx
from openpyxl import load_workbook # Para arquivos .xlsx

# Importações para Inteligência Artificial e Embeddings
import google.generativeai as genai
from sentence_transformers import SentenceTransformer

# Tentativa de importar FAISS, com fallback se não estiver disponível
try:
    import faiss
    FAISS_AVAILABLE = True
except ImportError:
    FAISS_AVAILABLE = False
    faiss = None # Define faiss como None para evitar erros em referências posteriores
import numpy as np

# Importações para a API web com Flask
from flask import Flask, request, jsonify, render_template
from flask_cors import CORS 
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address

# Configuração de ambiente (carrega variáveis do arquivo env)
from dotenv import load_dotenv
load_dotenv(dotenv_path='env_file.env') # Especifica o nome do seu arquivo env

# --- Configuração de Logging Melhorada ---
def setup_logging():
    log_level_str = os.getenv("LOG_LEVEL", "INFO").upper()
    log_file = os.getenv("LOG_FILE", "")
    log_format = "%(asctime)s - %(name)s - %(levelname)s - %(module)s:%(lineno)d - %(message)s"
    handlers = [logging.StreamHandler()]
    if log_file and log_file.strip():
        try:
            log_file_path = Path(log_file)
            log_file_path.parent.mkdir(parents=True, exist_ok=True)
            handlers.append(logging.FileHandler(log_file_path, encoding='utf-8'))
        except Exception as e:
            logging.error(f"Não foi possível configurar o logging para o arquivo {log_file}: {e}")
    log_level = getattr(logging, log_level_str, logging.INFO)
    logging.basicConfig(level=log_level, format=log_format, handlers=handlers)
    return logging.getLogger(__name__)

logger = setup_logging()

# --- Classes de Gerenciamento ---

class ConfigManager:
    def __init__(self):
        logger.debug("Inicializando ConfigManager...")
        self.google_api_key: str = os.getenv("GOOGLE_API_KEY", "")
        # Tentativa de conversão com tratamento de erro mais explícito
        try:
            self.chunk_size: int = int(os.getenv("CHUNK_SIZE", "300"))
            self.chunk_overlap: int = int(os.getenv("CHUNK_OVERLAP", "50"))
            self.max_context_docs: int = int(os.getenv("MAX_CONTEXT_DOCS", "5"))
            self.cache_ttl: int = int(os.getenv("CACHE_TTL_SECONDS", "86400")) 
            self.gemini_max_tokens: int = int(os.getenv("GEMINI_MAX_TOKENS", "2048")) 
            self.gemini_temperature: float = float(os.getenv("GEMINI_TEMPERATURE", "0.7"))
            self.flask_host: str = os.getenv("FLASK_HOST", "0.0.0.0") # LINHA ORIGINAL RESTAURADA
            # self.flask_host: str = "0.0.0.0" # LINHA DE TESTE TEMPORÁRIA REMOVIDA OU COMENTADA
            self.flask_port: int = int(os.getenv("FLASK_PORT", "5000"))
        except ValueError as e:
            logger.error(f"Erro ao converter valor do .env para número: {e}. Verifique os tipos no .env.")
            raise # Re-levanta a exceção para ser capturada por initialize_global_components

        self.embedding_model_name: str = os.getenv("EMBEDDING_MODEL", "intfloat/multilingual-e5-large")
        self.gemini_model_name: str = os.getenv("GEMINI_MODEL", "gemini-pro")
        self.data_dir: str = os.getenv("DATA_DIR", "dados_rag")
        self.target_document_name: str = "IDDC.pdf"
        self.cache_dir: str = os.getenv("CACHE_DIR", ".cache_embeddings") 
        self.enable_cache: bool = os.getenv("ENABLE_EMBEDDING_CACHE", "True").lower() == "true"
        
        self._validate_config()
        logger.info("ConfigManager inicializado e configurações validadas.")
    
    def _validate_config(self):
        logger.debug("Validando configurações...")
        if not self.google_api_key or self.google_api_key == "your_google_api_key_here":
            logger.warning("⚠️ GOOGLE_API_KEY não está configurada ou é um placeholder. Funcionalidades do Gemini podem ser limitadas.")
        if self.chunk_size <= 0: raise ValueError("CHUNK_SIZE deve ser maior que 0")
        if self.chunk_overlap < 0: raise ValueError("CHUNK_OVERLAP não pode ser negativo")
        if self.chunk_overlap >= self.chunk_size: raise ValueError("CHUNK_OVERLAP deve ser menor que CHUNK_SIZE")
        if self.max_context_docs <= 0: raise ValueError("MAX_CONTEXT_DOCS deve ser maior que 0")
        Path(self.data_dir).mkdir(parents=True, exist_ok=True) 
        logger.debug("Validação de configurações concluída.")

class CacheManager:
    def __init__(self, cache_dir: str, ttl: int = 86400):
        self.cache_dir = Path(cache_dir)
        self.ttl_seconds = ttl 
        try:
            self.cache_dir.mkdir(parents=True, exist_ok=True)
            logger.info(f"CacheManager inicializado. Diretório de cache: {self.cache_dir}")
        except Exception as e: logger.error(f"Erro ao criar diretório de cache {self.cache_dir}: {e}")
    
    def _get_cache_path(self, key: str) -> Optional[Path]:
        if not self.cache_dir: 
            logger.warning("CacheManager: diretório de cache não definido.")
            return None
        try:
            hash_key = hashlib.md5(key.encode('utf-8')).hexdigest()
            return self.cache_dir / f"{hash_key}.pkl"
        except Exception as e:
            logger.error(f"Erro ao gerar hash para a chave de cache: {e}")
            return None

    def get(self, key: str) -> Optional[np.ndarray]:
        cache_path = self._get_cache_path(key)
        if not cache_path or not cache_path.exists():
            logger.debug(f"Cache miss para a chave: {key[:50]}...")
            return None
        try:
            file_mod_time = cache_path.stat().st_mtime
            if (datetime.now().timestamp() - file_mod_time) > self.ttl_seconds:
                logger.info(f"Cache expirado para a chave: {key[:50]}... Removendo.")
                cache_path.unlink() 
                return None
            with open(cache_path, 'rb') as f: embedding = pickle.load(f)
            logger.debug(f"Cache hit para a chave: {key[:50]}...")
            return embedding
        except FileNotFoundError: 
             logger.debug(f"Cache miss (arquivo não encontrado após verificação) para a chave: {key[:50]}...")
             return None
        except Exception as e:
            logger.error(f"Erro ao carregar embedding do cache para a chave '{key[:50]}...': {e}")
            return None
    
    def set(self, key: str, value: np.ndarray):
        cache_path = self._get_cache_path(key)
        if not cache_path:
            logger.warning(f"Não foi possível salvar no cache para a chave: {key[:50]}... Caminho inválido.")
            return
        try:
            with open(cache_path, 'wb') as f: pickle.dump(value, f)
            logger.debug(f"Embedding salvo no cache para a chave: {key[:50]}...")
        except Exception as e: logger.error(f"Erro ao salvar embedding no cache para a chave '{key[:50]}...': {e}")

class DocumentProcessor:
    SUPPORTED_EXTENSIONS = {'.txt', '.pdf', '.pptx', '.csv', '.docx', '.xlsx'}
    @classmethod
    def is_supported(cls, file_path: Union[str, Path]) -> bool: return Path(file_path).suffix.lower() in cls.SUPPORTED_EXTENSIONS
    @staticmethod
    def read_txt(path: Union[str, Path]) -> str:
        logger.debug(f"Lendo arquivo TXT: {path}")
        try:
            with open(path, 'r', encoding='utf-8') as f: content = f.read()
            if not content.strip(): logger.warning(f"Arquivo TXT '{path}' está vazio (UTF-8).")
            return content
        except UnicodeDecodeError:
            logger.warning(f"Falha UTF-8 em '{path}'. Tentando Latin-1.")
            try:
                with open(path, 'r', encoding='latin-1') as f: content = f.read()
                if not content.strip(): logger.warning(f"Arquivo TXT '{path}' está vazio (Latin-1).")
                return content
            except Exception as e: logger.error(f"Erro Latin-1 TXT '{path}': {e}"); return ""
        except FileNotFoundError: logger.error(f"TXT não encontrado: {path}"); return ""
        except Exception as e: logger.error(f"Erro TXT '{path}': {e}"); return ""
    @staticmethod
    def read_pdf(path: Union[str, Path]) -> str:
        logger.debug(f"Lendo arquivo PDF: {path}")
        text = ""
        try:
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text(); text += (page_text + "\n") if page_text else ""
            if not text.strip(): logger.warning(f"PDF '{path}' vazio.")
            return text.strip()
        except Exception as e: logger.error(f"Erro PDF '{path}': {e}"); return ""
    @staticmethod
    def read_docx(path: Union[str, Path]) -> str:
        logger.debug(f"Lendo arquivo DOCX: {path}")
        try:
            doc = docx.Document(path); content = "\n".join([p.text for p in doc.paragraphs])
            if not content.strip(): logger.warning(f"DOCX '{path}' vazio.")
            return content.strip()
        except Exception as e: logger.error(f"Erro DOCX '{path}': {e}"); return ""
    @staticmethod
    def read_pptx(path: Union[str, Path]) -> str:
        logger.debug(f"Lendo arquivo PPTX: {path}")
        text = ""
        try:
            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
            if not text.strip(): logger.warning(f"PPTX '{path}' vazio.")
            return text.strip()
        except Exception as e: logger.error(f"Erro PPTX '{path}': {e}"); return ""
    @staticmethod
    def read_csv(path: Union[str, Path]) -> str:
        logger.debug(f"Lendo arquivo CSV: {path}")
        try:
            df = None
            try: df = pd.read_csv(path, encoding='utf-8')
            except UnicodeDecodeError: logger.warning(f"UTF-8 falhou CSV '{path}'. Tentando Latin-1."); df = pd.read_csv(path, encoding='latin-1')
            except pd.errors.ParserError:
                logger.warning(f"Parse falhou CSV '{path}' com vírgula. Tentando ';'.")
                try: df = pd.read_csv(path, sep=';', encoding='utf-8')
                except UnicodeDecodeError: df = pd.read_csv(path, sep=';', encoding='latin-1')
                except pd.errors.ParserError: logger.error(f"Parse CSV '{path}' falhou com ';'."); return ""
            if df is None: logger.error(f"Não carregou DataFrame CSV '{path}'."); return ""
            content = df.to_string(index=False)
            if not content.strip(): logger.warning(f"CSV '{path}' vazio.")
            return content
        except Exception as e: logger.error(f"Erro CSV '{path}': {e}"); return ""
    @staticmethod
    def read_xlsx(path: Union[str, Path]) -> str:
        logger.debug(f"Lendo arquivo XLSX: {path}")
        text_content = []
        try:
            wb = load_workbook(filename=path, read_only=True, data_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]; text_content.append(f"Planilha: {sheet_name}\n")
                for row in sheet.iter_rows():
                    text_content.append(", ".join([str(c.value) if c.value is not None else "" for c in row]) + "\n")
                text_content.append("\n")
            content = "".join(text_content)
            if not content.strip(): logger.warning(f"XLSX '{path}' vazio.")
            return content.strip()
        except Exception as e: logger.error(f"Erro XLSX '{path}': {e}"); return ""
    @classmethod
    def extract_text(cls, file_path: Union[str, Path]) -> str:
        path_obj = Path(file_path); ext = path_obj.suffix.lower()
        logger.info(f"Extraindo texto de: {file_path} (ext: {ext})")
        if not cls.is_supported(path_obj): logger.warning(f"Tipo não suportado: {ext} para {file_path}"); return ""
        try:
            if ext == '.txt': return cls.read_txt(path_obj)
            elif ext == '.pdf': return cls.read_pdf(path_obj)
            elif ext == '.docx': return cls.read_docx(path_obj)
            elif ext == '.pptx': return cls.read_pptx(path_obj)
            elif ext == '.csv': return cls.read_csv(path_obj)
            elif ext == '.xlsx': return cls.read_xlsx(path_obj)
            else: logger.error(f"Lógica não implementada para ext: {ext}"); return ""
        except Exception as e: logger.error(f"Erro extraindo de {file_path}: {e}", exc_info=True); return ""

class RAGSystem:
    def __init__(self, config_manager: ConfigManager, embedding_model_name: str, gemini_model_name: str):
        self.config = config_manager
        self.logger = logging.getLogger(__name__ + ".RAGSystem") 
        self.embedding_model: Optional[SentenceTransformer] = None
        self.gemini_model: Optional[genai.GenerativeModel] = None
        self.vector_store = None 
        self.documents_processed = False 
        self.chunk_texts: List[str] = [] 
        self.cache_manager: Optional[CacheManager] = None

        if self.config.enable_cache:
            self.cache_manager = CacheManager(cache_dir=self.config.cache_dir, ttl=self.config.cache_ttl)
            self.logger.info(f"Cache de embeddings habilitado. Diretório: {self.config.cache_dir}, TTL: {self.config.cache_ttl}s")
        else:
            self.logger.info("Cache de embeddings desabilitado.")
        self.last_retrieved_contexts_for_eval: List[str] = [] # Para RAGAS
        try:
            self.logger.info(f"Carregando modelo de embedding: {embedding_model_name}")
            self.embedding_model = SentenceTransformer(embedding_model_name)
            self.logger.info("Modelo de embedding carregado.")
        except Exception as e: self.logger.error(f"Falha ao carregar embedding model: {e}", exc_info=True)
        if self.config.google_api_key and self.config.google_api_key != "your_google_api_key_here":
            try:
                self.logger.info(f"Configurando modelo Gemini: {gemini_model_name}")
                self.gemini_model = genai.GenerativeModel(gemini_model_name)
                self.logger.info("Modelo Gemini configurado.")
            except Exception as e: self.logger.error(f"Falha ao configurar Gemini: {e}", exc_info=True)
        else: self.logger.warning("API Key Gemini não configurada. Modelo Gemini não inicializado.")
        self._initialize_vector_store()
        # Processa o documento alvo automaticamente na inicialização do RAGSystem
        self.process_and_index_documents()

    def _initialize_vector_store(self):
        self.logger.info("Inicializando vector store...")
        if FAISS_AVAILABLE:
            embedding_dim = 1024 
            try:
                self.vector_store = faiss.IndexFlatL2(embedding_dim)
                self.logger.info(f"Índice FAISS (IndexFlatL2) criado (dim {embedding_dim}).")
            except Exception as e: self.logger.error(f"Falha ao criar índice FAISS: {e}", exc_info=True); self.vector_store = None 
        else:
            self.logger.warning("FAISS não disponível. Usando lista fallback (limitado).")
            self.vector_store = [] 

    def _get_text_chunks(self, text: str) -> List[str]:
        if not text or not text.strip(): return []
        chunks = []; start = 0; text_len = len(text)
        while start < text_len:
            end = start + self.config.chunk_size; chunk = text[start:end]; chunks.append(chunk)
            if end >= text_len: break
            start += (self.config.chunk_size - self.config.chunk_overlap)
            if start >= text_len and end < text_len : break
        return chunks

    def process_and_index_documents(self):
        if not self.embedding_model: self.logger.error("Embedding model não carregado."); return
        doc_dir = Path(self.config.data_dir); target_file_path = doc_dir / self.config.target_document_name
        self.logger.info(f"Processando documento alvo: {target_file_path}")
        self._initialize_vector_store(); self.chunk_texts = []; self.documents_processed = False
        if not target_file_path.is_file(): self.logger.error(f"'{self.config.target_document_name}' não encontrado em '{doc_dir}'."); return
        if not DocumentProcessor.is_supported(target_file_path): self.logger.error(f"'{self.config.target_document_name}' tipo não suportado."); return
        
        texts_this_run, embeds_this_run = [], []
        
        text_content = DocumentProcessor.extract_text(target_file_path)
        if text_content:
            chunks = self._get_text_chunks(text_content)
            if chunks:
                self.logger.info(f"Processando {len(chunks)} chunks de '{target_file_path.name}'. Verificando cache...")
                
                cached_embeddings_map: Dict[str, np.ndarray] = {}
                chunks_to_encode_texts: List[str] = []

                for i, chunk_text in enumerate(chunks):
                    cached_embedding = None
                    if self.cache_manager and self.config.enable_cache:
                        cached_embedding = self.cache_manager.get(chunk_text)
                    
                    if cached_embedding is not None:
                        cached_embeddings_map[chunk_text] = cached_embedding
                        # self.logger.debug(f"Cache hit para chunk {i+1}/{len(chunks)}.")
                    else:
                        chunks_to_encode_texts.append(chunk_text)
                        # self.logger.debug(f"Cache miss para chunk {i+1}/{len(chunks)}. Agendado para encoding.")

                if chunks_to_encode_texts:
                    self.logger.info(f"Gerando embeddings para {len(chunks_to_encode_texts)} chunks não cacheados...")
                    try:
                        newly_generated_embeddings = self.embedding_model.encode(
                            chunks_to_encode_texts, 
                            convert_to_tensor=False, 
                            show_progress_bar=False # Mantenha False para evitar output excessivo
                        )
                        for i, text_chunk_encoded in enumerate(chunks_to_encode_texts):
                            embedding_value = newly_generated_embeddings[i]
                            cached_embeddings_map[text_chunk_encoded] = embedding_value # Adiciona ao mapa geral
                            if self.cache_manager and self.config.enable_cache:
                                self.cache_manager.set(text_chunk_encoded, embedding_value)
                        self.logger.info(f"Gerados e cacheados {len(newly_generated_embeddings)} novos embeddings.")
                    except Exception as e:
                        self.logger.error(f"Erro ao gerar/cachear novos embeddings para {target_file_path.name}: {e}", exc_info=True)
                
                # Reconstruir texts_this_run e embeds_this_run na ordem original dos chunks
                for chunk_text in chunks:
                    if chunk_text in cached_embeddings_map:
                        texts_this_run.append(chunk_text)
                        embeds_this_run.append(cached_embeddings_map[chunk_text])

            else: self.logger.warning(f"Nenhum chunk gerado para {target_file_path.name}")
        else: self.logger.warning(f"Nenhum texto extraído de {target_file_path.name}")

        if embeds_this_run:
            embeddings_np = np.array(embeds_this_run).astype('float32') 
            if FAISS_AVAILABLE and self.vector_store is not None and hasattr(self.vector_store, 'add'):
                self.vector_store.add(embeddings_np)
                self.logger.info(f"Adicionados {len(embeds_this_run)} vetores de '{self.config.target_document_name}' ao FAISS. Total no índice: {self.vector_store.ntotal}.")
            elif isinstance(self.vector_store, list): 
                for emb, txt in zip(embeddings_np, texts_this_run): self.vector_store.append((emb, txt))
                self.logger.info(f"Adicionados {len(self.vector_store)} embeds/textos de '{self.config.target_document_name}' ao fallback.")
            self.chunk_texts = texts_this_run # Atualiza a lista de textos na ordem correta
            self.documents_processed = True
        else: self.logger.warning(f"Nenhum embedding gerado/adicionado para '{self.config.target_document_name}'.")

    def get_relevant_context(self, question_embedding: np.ndarray, k: int) -> List[str]:
        if not self.documents_processed: self.logger.warning(f"'{self.config.target_document_name}' não processado."); return []
        if self.vector_store is None: self.logger.warning("Vector store não inicializado."); return []
        if FAISS_AVAILABLE and hasattr(self.vector_store, 'ntotal') and self.vector_store.ntotal == 0: self.logger.warning("FAISS vazio."); return []
        if isinstance(self.vector_store, list) and not self.vector_store: self.logger.warning("Fallback store vazio."); return []
        
        relevant_texts = []
        if FAISS_AVAILABLE and hasattr(self.vector_store, 'search') and self.vector_store.ntotal > 0 :
            try:
                actual_k = min(k, self.vector_store.ntotal)
                if actual_k == 0: return [] 
                _, indices = self.vector_store.search(np.array([question_embedding]).astype('float32'), actual_k)
                self.logger.debug(f"Índices FAISS recuperados: {indices[0]}")
                for i in indices[0]: 
                    if 0 <= i < len(self.chunk_texts): relevant_texts.append(self.chunk_texts[i])
                    else: self.logger.warning(f"Índice FAISS {i} fora de alcance (len: {len(self.chunk_texts)})")
                self.logger.debug(f"Contexto FAISS (índices): {indices[0]}")
            except Exception as e: self.logger.error(f"Erro busca FAISS: {e}", exc_info=True); return [] 
        elif isinstance(self.vector_store, list) and self.vector_store: 
            self.logger.warning("Usando busca similaridade cosseno (ineficiente).")
            similarities = []
            for stored_emb, text_chunk in self.vector_store:
                norm_q = question_embedding / (np.linalg.norm(question_embedding) + 1e-8) 
                norm_s = stored_emb / (np.linalg.norm(stored_emb) + 1e-8)
                similarities.append((np.dot(norm_q, norm_s), text_chunk))
            similarities.sort(key=lambda x: x[0], reverse=True)
            relevant_texts = [text for _, text in similarities[:k]]
            self.logger.debug(f"Contexto fallback (top {k} sims): {[s[0] for s in similarities[:k]]}")
        self.last_retrieved_contexts_for_eval = relevant_texts # Armazena para RAGAS
        return relevant_texts

    def get_answer(self, question: str) -> str:
        if not self.embedding_model: return "Erro: Embedding model não pronto."
        if not self.gemini_model: return "Erro: Gemini model não pronto."
        if not self.documents_processed: return f"'{self.config.target_document_name}' precisa ser processado. Use 'rag-gemini process'."
        try:
            self.logger.info(f"Gerando embedding para pergunta: '{question[:50]}...'")
            q_embed = self.embedding_model.encode(question, convert_to_tensor=False, show_progress_bar=False)
        except Exception as e: self.logger.error(f"Erro embedding pergunta: {e}", exc_info=True); return "Erro processando pergunta (embedding)."
        
        self.logger.info(f"Buscando contexto de '{self.config.target_document_name}'...")
        ctx_chunks = self.get_relevant_context(q_embed, k=self.config.max_context_docs)
        ctx_text = f"\n\n---\n\n".join(ctx_chunks) if ctx_chunks else f"Nenhum contexto de '{self.config.target_document_name}'."
        if not ctx_chunks: self.logger.warning(f"Nenhum contexto de '{self.config.target_document_name}' para pergunta.")
        else: self.logger.debug(f"Contexto LLM (primeiros 200 chars): {ctx_text[:200]}...")

        prompt = f"""Você é um assistente especializado e prestativo. Sua principal tarefa é responder perguntas com base no conteúdo do documento '{self.config.target_document_name}'.
Utilize o contexto fornecido abaixo, extraído deste documento, como a fonte primária para sua resposta.
Você pode sintetizar e elaborar a informação encontrada no contexto para fornecer uma resposta clara e coesa.
Evite usar conhecimento externo ao documento.
Se a informação necessária para responder à pergunta não puder ser razoavelmente inferida ou encontrada no contexto fornecido, indique que a informação específica não está detalhada no documento '{self.config.target_document_name}'.

Contexto do documento {self.config.target_document_name}:
{ctx_text}

Pergunta do Usuário:
{question}

Resposta (elaborada com base no contexto do documento '{self.config.target_document_name}'):
"""
        self.logger.info("Enviando prompt para Gemini..."); self.logger.debug(f"Prompt (500 chars):\n{prompt[:500]}")
        try:
            gen_config = genai.types.GenerationConfig(candidate_count=1, max_output_tokens=self.config.gemini_max_tokens, temperature=self.config.gemini_temperature)
            safety = [{"category": c, "threshold": "BLOCK_NONE"} for c in ["HARM_CATEGORY_HARASSMENT", "HARM_CATEGORY_HATE_SPEECH", "HARM_CATEGORY_SEXUALLY_EXPLICIT", "HARM_CATEGORY_DANGEROUS_CONTENT"]]
            response = self.gemini_model.generate_content(prompt, generation_config=gen_config, safety_settings=safety)
            
            if response.parts: answer_text = "".join(p.text for p in response.parts if hasattr(p, 'text'))
            elif response.candidates and response.candidates[0].content and response.candidates[0].content.parts:
                answer_text = "".join(p.text for p in response.candidates[0].content.parts if hasattr(p, 'text'))
            elif response.prompt_feedback and response.prompt_feedback.block_reason:
                reason = response.prompt_feedback.block_reason.name
                msg = f" ({response.prompt_feedback.block_reason_message})" if hasattr(response.prompt_feedback, 'block_reason_message') and response.prompt_feedback.block_reason_message else ""
                self.logger.warning(f"Geração bloqueada. Razão: {reason}{msg}")
                return f"Desculpe, resposta bloqueada por restrições de conteúdo ({reason})."
            else: self.logger.warning(f"Resposta Gemini sem partes/bloqueada sem razão clara. Detalhes: {response}"); return "Desculpe, não obtive resposta clara."
            self.logger.info(f"Resposta Gemini: {answer_text[:100]}..."); return answer_text.strip()
        except Exception as e: self.logger.error(f"Erro API Gemini: {e}", exc_info=True); return f"Erro comunicação LLM: {str(e)}"

    def get_last_retrieved_contexts(self) -> List[str]:
        """Retorna os últimos contextos recuperados para avaliação com RAGAS."""
        return self.last_retrieved_contexts_for_eval

app = Flask(__name__)
CORS(app) 

config: Optional[ConfigManager] = None
rag_system_instance: Optional[RAGSystem] = None

def initialize_global_components():
    global config, rag_system_instance
    if config is not None and rag_system_instance is not None: # Evita reinicialização
        logger.debug("Componentes globais já inicializados.")
        return
    try:
        logger.info("Inicializando componentes globais...")
        current_config = ConfigManager()
        logger.info("ConfigManager carregado e validado.")
        if current_config.google_api_key and current_config.google_api_key != "your_google_api_key_here":
            genai.configure(api_key=current_config.google_api_key)
            logger.info("Google Generative AI configurado.")
        else: logger.error("GOOGLE_API_KEY não configurada. Gemini indisponível.")
        
        current_rag_system = RAGSystem(
            config_manager=current_config,
            embedding_model_name=current_config.embedding_model_name,
            gemini_model_name=current_config.gemini_model_name
        )
        logger.info("Instância RAGSystem criada.")
        if not FAISS_AVAILABLE: logger.warning(f"FAISS não instalado. Funcionalidade para '{current_config.target_document_name}' limitada.")
        logger.info(f"Componentes RAG inicializados, foco em '{current_config.target_document_name}'.")
        
        config = current_config
        rag_system_instance = current_rag_system
    except ValueError as ve:
        logger.critical(f"Erro de VALOR na configuração durante a inicialização global: {ve}. Verifique seu arquivo .env.")
        logger.critical("A APLICAÇÃO NÃO PODE CONTINUAR SEM CONFIGURAÇÃO VÁLIDA. ENCERRANDO.")
        sys.exit(1) # Sai imediatamente se ConfigManager falhar com ValueError
    except Exception as e:
        logger.critical(f"Erro fatal durante a inicialização global: {e}", exc_info=True)
        logger.critical("A APLICAÇÃO NÃO PODE CONTINUAR DEVIDO A UM ERRO FATAL NA INICIALIZAÇÃO. ENCERRANDO.")
        sys.exit(1) # Sai imediatamente para outros erros fatais durante a inicialização

initialize_global_components() # Chama a inicialização quando o módulo é carregado

@app.route('/')
def rota_principal_health_check():
    logger.debug("Endpoint '/' acessado.")
    doc_name = config.target_document_name if config else "documento alvo"
    return jsonify({"status": "Online!", "timestamp": datetime.now().isoformat(), "message": f"API RAG Gemini (foco: {doc_name})."})

@app.route('/api/ask', methods=['POST'])
def rota_ask_question_api():
    global rag_system_instance 
    try:
        data = request.get_json()
        if not data or 'question' not in data: return jsonify({"error": "'question' obrigatório."}), 400
        pergunta = data['question']
        if not isinstance(pergunta, str) or not pergunta.strip(): return jsonify({"error": "'question' deve ser string não vazia."}), 400
        logger.info(f"API /api/ask: Pergunta: '{pergunta[:100]}...'") 
        if rag_system_instance is None: return jsonify({"error": "RAG não inicializado."}), 503
        if not rag_system_instance.documents_processed:
            return jsonify({"error": f"'{rag_system_instance.config.target_document_name}' precisa ser processado. Use 'rag-gemini process'."}), 400
        resposta = rag_system_instance.get_answer(pergunta)
        logger.info(f"API /api/ask: Resposta: '{resposta[:100]}...'")
        return jsonify({"question": pergunta, "answer": resposta})
    except json.JSONDecodeError: return jsonify({"error": "JSON inválido."}), 400
    except Exception as e: logger.error(f"API /api/ask erro: {e}", exc_info=True); return jsonify({"error": "Erro interno."}), 500

@app.route('/ui', methods=['GET'])
def interface_usuario_web():
    logger.debug("Acessando /ui")
    try:
        doc_alvo = config.target_document_name if config else "documento configurado"
        return render_template('index.html', title="ChatBot COD")
    except Exception as e:
        logger.error(f"Erro renderizando template para /ui: {e}", exc_info=True)
        return "Erro ao carregar interface. Verifique logs.", 500

def processar_documentos_cli():
    logger.info(f"CLI: Processando '{config.target_document_name if config else 'documento alvo'}'...")
    if rag_system_instance:
        rag_system_instance.process_and_index_documents()
        logger.info(f"CLI: Processamento de '{config.target_document_name if config else 'documento alvo'}' concluído.")
    else: logger.error("CLI 'process': RAGSystem não disponível.")

def fazer_pergunta_cli(texto_pergunta: str):
    logger.info(f"CLI: Pergunta: '{texto_pergunta[:100]}...'")
    if rag_system_instance:
        if not rag_system_instance.documents_processed:
            print(f"\nAVISO: '{rag_system_instance.config.target_document_name}' não processado. Execute 'rag-gemini process'.")
        resposta = rag_system_instance.get_answer(texto_pergunta)
        print(f"\nResposta RAG ('{rag_system_instance.config.target_document_name}'):\n{resposta}")
        logger.info("CLI: Pergunta respondida.")
    else: print("\nErro: RAGSystem não disponível."); logger.error("CLI 'ask': RAGSystem não disponível.")

def run_server():
    if config is None: logger.critical("Configuração não carregada. Servidor não pode iniciar."); return
    if not config.google_api_key or config.google_api_key == "your_google_api_key_here":
        logger.error("GOOGLE_API_KEY não configurada. Servidor pode não funcionar como esperado.")
    host, port = config.flask_host, config.flask_port
    debug = os.getenv("FLASK_DEBUG", "False").lower() == "true" # FLASK_DEBUG pode continuar sendo lido diretamente ou movido para ConfigManager também
    logger.info(f"🚀 Iniciando Flask em http://{host}:{port} (Debug: {debug})")
    try: app.run(host=host, port=port, debug=debug, use_reloader=debug)
    except Exception as e: logger.critical(f"Falha Flask: {e}", exc_info=True)

def main():
    global rag_system_instance, config
    # A inicialização global já foi chamada. Verificamos aqui se foi bem-sucedida.
    if config is None: 
        logger.critical("Falha crítica na inicialização dos componentes globais (config é None). Encerrando.")
        # Não há necessidade de chamar initialize_global_components() novamente aqui,
        # pois ela já foi chamada e, se falhou, já logou e saiu (com a nova lógica de sys.exit).
        # Se chegou aqui e config é None, é um estado inesperado ou a lógica de saída falhou.
        sys.exit(1)


    doc_alvo = config.target_document_name if hasattr(config, 'target_document_name') else "documento alvo"
    parser = argparse.ArgumentParser(description=f"RAG Gemini (foco: {doc_alvo}).", formatter_class=argparse.RawTextHelpFormatter)
    subparsers = parser.add_subparsers(dest="command", help="Comandos:")
    subparsers.add_parser("server", help="Iniciar servidor web Flask.")
    subparsers.add_parser("process", help=f"Processar '{doc_alvo}' e construir índice.")
    ask_p = subparsers.add_parser("ask", help="Fazer pergunta via CLI.")
    ask_p.add_argument("texto_pergunta", metavar="PERGUNTA", type=str, help="Texto da pergunta.")

    # Modificado para lidar com o caso de nenhum argumento ser passado para 'rag-gemini'
    if len(sys.argv) <= 1: # Se apenas 'rag-gemini' (ou o nome do script) for executado
        args = argparse.Namespace(command="server") # Default para o comando 'server'
        logger.info("Nenhum comando CLI especificado. Iniciando servidor por padrão...")
    else:
        args = parser.parse_args(sys.argv[1:]) # Processa os argumentos normalmente
    
    # Se, mesmo após o parsing, args.command for None (o que não deveria acontecer com o default acima)
    if not hasattr(args, 'command') or args.command is None:
        args.command = "server" # Garante um comando default
        logger.info("Comando não reconhecido ou ausente. Iniciando servidor por padrão...")


    logger.info(f"🚀 RAG Gemini via main() comando: {args.command}")
    if args.command == "server": run_server()
    elif args.command == "process": processar_documentos_cli()
    elif args.command == "ask": fazer_pergunta_cli(args.texto_pergunta)

if __name__ == '__main__':
    main()